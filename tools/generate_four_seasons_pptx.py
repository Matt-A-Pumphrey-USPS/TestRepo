"""
generate_four_seasons_pptx.py

Generates a 5-slide PowerPoint (Title + Spring/Summer/Autumn/Winter) with light humor and speaker notes.
Requires: python3 and python-pptx (pip install python-pptx)

Usage:
  python generate_four_seasons_pptx.py

This will create Four_Seasons_generated.pptx in the current directory.
"""
from pptx import Presentation
from pptx.util import Inches, Pt

slides_content = [
    {
        "title": "Four Seasons",
        "subtitle": "A weather report with attitude",
        "notes": "Welcome — keep it short and smile."
    },
    {
        "title": "Spring",
        "bullets": [
            "Flowers come back (and so does the yard work)",
            "Allergies throw a surprise party",
            "Time for optimism (and muddy shoes)"
        ],
        "notes": "Make a joke about pollen being nature's confetti."
    },
    {
        "title": "Summer",
        "bullets": [
            "Sunscreen becomes a social contract",
            "Ice cream sales spike (science)",
            "Long days, questionable decisions around 9pm"
        ],
        "notes": "Mention the unofficial summer sport: evaporating ice cubes."
    },
    {
        "title": "Autumn",
        "bullets": [
            "Leaves perform their dramatic exits",
            "Coffee gets cozier",
            "Sweaters stage a comeback tour"
        ],
        "notes": "Call out the leaf-blower vs. leaf-pile debate."
    },
    {
        "title": "Winter",
        "bullets": [
            "Snow: magical until you're shoveling it",
            "Holiday lights: engineering challenges",
            "Hot beverages reach peak popularity"
        ],
        "notes": "Remind everyone that winter is nature's nap time."
    }
]

prs = Presentation()
# Use built-in title slide layout for slide 1
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = slides_content[0]["title"]
subtitle.text = slides_content[0]["subtitle"]
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = slides_content[0]["notes"]

# For subsequent slides use a title and content layout (usually layout 1)
bullet_layout = prs.slide_layouts[1]
for s in slides_content[1:]:
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = s["title"]
    # find the first placeholder for content and add bullets
    body = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # BODY
            body = shape
            break
    if body is None:
        # fallback: add a textbox
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)
        height = Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
    else:
        tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(s["bullets"]):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
    # speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = s["notes"]

output_filename = "Four_Seasons_generated.pptx"
prs.save(output_filename)
print(f"Generated PPTX: {output_filename}")
